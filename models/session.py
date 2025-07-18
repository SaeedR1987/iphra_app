from docxtpl import DocxTemplate
from pptx import Presentation
import pandas as pd
from datetime import datetime
from typing import Optional, Union
import pandas as pd
import statsmodels.api as sm
import numpy as np

class Session:
    """
    A class to represent a session in the application.
    This class is designed to manage the session state, including user authentication,
    user roles, and session data.
    """

    SUPPORTED_VERSION = "1.0.0"
    APP_VERSION = "1.0.0"
    APP_NAME = "IPHRA_APP"

    created_at = datetime.now().isoformat()
    modified_at = None
    last_saved_at = None
    has_been_validated = False
    is_complete = False



    def __init__(self):
        self.sample = Sample()
        self.tool_household = HouseholdTool()
        self.tool_community_kii = CommunityKiiTool()
        self.tool_health_facility_kii = HealthFacilityKiiTool()
        self.tool_health_facility_obs = HealthFacilityObservationTool()
        self.tool_fsl_kii = FslKiiTool()
        self.tool_nut_kii = NutKiiTool()
        self.tool_water_obs = WaterObservationTool()
        self.tool_latrine_obs = LatrineObservationTool()
        self.data = Data()
        self.cleaning_log = []
        self.quality_report = {}
        self.analysis = Analysis()
        self.graphics_log = []

    def update_modified_time(self):
        self.metadata["modified_at"] = datetime.now().isoformat()

    def to_dict(self):
        """
        Convert the session data to a dictionary format.
        This method should return the session data in a format suitable for serialization or storage.
        """
        return {
            "sample": self.sample.to_dict(),
            "tool_household": self.tool_household.to_dict(),
            "tool_community_kii": self.tool_community_kii.to_dict(),
            "tool_health_facility_kii": self.tool_health_facility_kii.to_dict(),
            "tool_health_facility_obs": self.tool_health_facility_obs.to_dict(),
            "tool_fsl_kii": self.tool_fsl_kii.to_dict(),
            "tool_nut_kii": self.tool_nut_kii.to_dict(),
            "tool_water_obs": self.tool_water_obs.to_dict(),
            "tool_latrine_obs": self.tool_latrine_obs.to_dict()
        }
    
    @staticmethod
    def from_dict(self, data):
        """
        Load the session data from a dictionary format.
        This method should populate the session attributes from the provided dictionary.
        """
        self.sample.from_dict(data.get("sample", {}))
        self.tool_household.from_dict(data.get("tool_household", {}))
        self.tool_community_kii.from_dict(data.get("tool_community_kii", {}))
        self.tool_health_facility_kii.from_dict(data.get("tool_health_facility_kii", {}))
        self.tool_health_facility_obs.from_dict(data.get("tool_health_facility_obs", {}))
        self.tool_fsl_kii.from_dict(data.get("tool_fsl_kii", {}))
        self.tool_nut_kii.from_dict(data.get("tool_nut_kii", {}))
        self.tool_water_obs.from_dict(data.get("tool_water_obs", {}))
        self.tool_latrine_obs.from_dict(data.get("tool_latrine_obs", {}))

    def save_to_file(self, path: str):
        import json
        with open(path, "w") as f:
            json.dump(self.to_dict(), f, indent=4)

    @staticmethod
    def load_from_file(path: str):
        import json
        with open(path, "r") as f:
            data = json.load(f)
        return Session.from_dict(data)
    
    def reset(self):
        self.sample = Sample()
        self.tool = Tool()
        self.analysis = Analysis()
        self.notes = ""
        self.metadata = {}

    def set_metadata(self, key, value):
        self.metadata[key] = value

    def get_metadata(self, key, default=None):
        return self.metadata.get(key, default)

    def is_valid(self):
        """
        Check if the session is valid.
        This method should verify that all necessary components of the session are initialized.
        """
        return (
            self.sample is not None and
            self.tool is not None and
            self.analysis is not None
        )

    def generate_integrated_analysis
        pass
    

    def generate_word_report(session, template_path, output_path):
        doc = DocxTemplate(template_path)
        
        context = {
            "sample_design": session.sample.design,
            "mortality_rate": session.sample.mortality_rate,
            "notes": session.notes,
            # Add all fields you want in your report
        }
        
        doc.render(context)
        doc.save(output_path)

    def generate_ppt_report(session, template_path, output_path):
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "${sample_design}" in text:
                        shape.text_frame.text = text.replace("${sample_design}", session.sample.design)
                    # Repeat for other placeholders
                    
        prs.save(output_path)
